from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Brand Colors ──────────────────────────────────────────────────────────────
BURGUNDY   = RGBColor(0x6B, 0x2D, 0x39)
DARK_BURG  = RGBColor(0x4A, 0x1F, 0x28)
MED_BURG   = RGBColor(0x8B, 0x4D, 0x5A)
GOLD       = RGBColor(0xD4, 0xA5, 0x73)
LIGHT_GOLD = RGBColor(0xED, 0xD9, 0xB8)
CREAM      = RGBColor(0xF8, 0xF7, 0xF5)
CREAM_WARM = RGBColor(0xFE, 0xF6, 0xEE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x0A, 0x0F)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

W = Inches(13.333)
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def box(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size, color,
        bold=False, italic=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.color.rgb = color
    r.font.bold   = bold
    r.font.italic = italic
    r.font.name   = font
    return tb


def txt_multi(slide, lines, l, t, w, h, align=PP_ALIGN.LEFT, font="Calibri"):
    """lines = list of (text, size, color, bold)"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, color, bold) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font
    return tb


def header(slide, title):
    """Standard burgundy header bar used on all content slides."""
    box(slide, Inches(0), Inches(0), W, Inches(1.15), BURGUNDY)
    box(slide, Inches(0), Inches(1.15), W, Inches(0.07), GOLD)
    txt(slide, title,
        Inches(0.45), Inches(0.2), Inches(12.4), Inches(0.8),
        26, WHITE, bold=True, font="Georgia")


def diagram_placeholder(slide, label):
    """Full-area placeholder for a Mermaid diagram PNG."""
    box(slide, Inches(0.5), Inches(1.38), Inches(12.33), Inches(5.87),
        CREAM_WARM, GOLD)
    # Corner accents
    for lx, ty in [(0.5, 1.38), (12.33, 1.38), (0.5, 6.75), (12.33, 6.75)]:
        box(slide, Inches(lx if lx < 1 else lx + 0.33),
            Inches(ty if ty < 2 else ty + 0.25),
            Inches(0.3), Inches(0.08), BURGUNDY)
    txt_multi(slide, [
        (f"[ Insert  {label}  PNG here ]", 18, MED_BURG, True),
        ("", 8, BLACK, False),
        ("How to insert:", 12, BLACK, True),
        ("1.  Open  mermaid.live  and paste the diagram code", 12, BLACK, False),
        ("2.  Click  Download PNG  (top-right of the preview panel)", 12, BLACK, False),
        ("3.  Delete this text box", 12, BLACK, False),
        ("4.  Click Insert  Picture  and choose the downloaded PNG", 12, BLACK, False),
        ("5.  Resize the image to fill this bordered area", 12, BLACK, False),
    ],
        Inches(3.8), Inches(2.8), Inches(6.2), Inches(3.5),
        align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
bg(s1, BURGUNDY)

# Gold top & bottom accent bars
box(s1, Inches(0), Inches(0),       W, Inches(0.18), GOLD)
box(s1, Inches(0), Inches(7.32),    W, Inches(0.18), GOLD)

# Main title
txt(s1, "EventWearify",
    Inches(0), Inches(1.6), W, Inches(1.6),
    68, WHITE, bold=True, align=PP_ALIGN.CENTER, font="Georgia")

# Gold divider
box(s1, Inches(4.2), Inches(3.35), Inches(4.9), Inches(0.07), GOLD)

# Subtitle
txt(s1, "Event Wear Lease System",
    Inches(0), Inches(3.5), W, Inches(0.75),
    26, GOLD, align=PP_ALIGN.CENTER, font="Georgia")

# Name
txt(s1, "Jasmine Ciely Palicte",
    Inches(0), Inches(4.5), W, Inches(0.55),
    20, WHITE, bold=True, align=PP_ALIGN.CENTER)

# Course
txt(s1, "IT342  —  Systems Integration and Architecture",
    Inches(0), Inches(5.1), W, Inches(0.5),
    15, LIGHT_GOLD, align=PP_ALIGN.CENTER)

# Year
txt(s1, "2nd Semester  ·  S.Y. 2025 – 2026",
    Inches(0), Inches(5.65), W, Inches(0.45),
    13, LIGHT_GOLD, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — System Introduction
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
bg(s2, CREAM)
header(s2, "SYSTEM PROJECT INTRODUCTION")

BW = Inches(3.9)
BH = Inches(4.8)
BY = Inches(1.5)

# Box 1 — Purpose (Burgundy)
box(s2, Inches(0.35), BY, BW, BH, BURGUNDY, DARK_BURG)
txt(s2, "PURPOSE", Inches(0.35), Inches(1.6), BW, Inches(0.55),
    17, GOLD, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
box(s2, Inches(0.55), Inches(2.22), Inches(3.5), Inches(0.05), GOLD)
txt_multi(s2, [
    ("A digital rental and lease management", 13, WHITE, False),
    ("platform for event wear shops —", 13, WHITE, False),
    ("businesses that rent out gowns, suits,", 13, WHITE, False),
    ("barongs, and special-occasion attire.", 13, WHITE, False),
    ("", 8, WHITE, False),
    ("EventWearify replaces manual, paper-based", 13, WHITE, False),
    ("booking methods with a fully digital,", 13, WHITE, False),
    ("centralized system accessible on both", 13, WHITE, False),
    ("web and Android mobile.", 13, WHITE, False),
], Inches(0.55), Inches(2.35), Inches(3.5), Inches(3.5))

# Box 2 — Problem (Gold)
box(s2, Inches(4.72), BY, BW, BH, GOLD, DARK_BURG)
txt(s2, "PROBLEM", Inches(4.72), Inches(1.6), BW, Inches(0.55),
    17, DARK_BURG, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
box(s2, Inches(4.92), Inches(2.22), Inches(3.5), Inches(0.05), BURGUNDY)
txt_multi(s2, [
    ("Most small event wear shops manage", 13, BLACK, False),
    ("bookings through text messages,", 13, BLACK, False),
    ("paper records, or spreadsheets.", 13, BLACK, False),
    ("", 8, BLACK, False),
    ("This leads to:", 13, BLACK, True),
    ("·  Double bookings", 13, BLACK, False),
    ("·  Missed reservations", 13, BLACK, False),
    ("·  No real-time visibility", 13, BLACK, False),
    ("·  Poor customer experience", 13, BLACK, False),
], Inches(4.92), Inches(2.35), Inches(3.5), Inches(3.5))

# Box 3 — Users (Cream Warm)
box(s2, Inches(9.09), BY, BW, BH, CREAM_WARM, BURGUNDY)
txt(s2, "USERS", Inches(9.09), Inches(1.6), BW, Inches(0.55),
    17, BURGUNDY, bold=True, align=PP_ALIGN.CENTER, font="Georgia")
box(s2, Inches(9.29), Inches(2.22), Inches(3.5), Inches(0.05), GOLD)
txt_multi(s2, [
    ("ADMIN", 14, BURGUNDY, True),
    ("Shop owner or staff who manage", 13, BLACK, False),
    ("inventory, bookings, customers,", 13, BLACK, False),
    ("and overall operations.", 13, BLACK, False),
    ("", 8, BLACK, False),
    ("CUSTOMER", 14, BURGUNDY, True),
    ("Clients who browse the catalog,", 13, BLACK, False),
    ("check availability, and place", 13, BLACK, False),
    ("bookings online.", 13, BLACK, False),
], Inches(9.29), Inches(2.35), Inches(3.5), Inches(3.7))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Main Features
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
bg(s3, CREAM)
header(s3, "MAIN FEATURES")

features = [
    ("01", "Authentication & Authorization",
     "JWT login with role-based access control (Admin / Customer) and Google OAuth2 social login."),
    ("02", "Inventory Management",
     "Add, update, and delete items with images stored on Supabase Cloud Storage."),
    ("03", "Fitting Bookings",
     "Schedule physical fitting appointments with time-slot availability enforcement."),
    ("04", "Direct Lease Bookings",
     "Book items for a date range with automatic price and weekly/monthly discount calculation."),
    ("05", "Custom Availability Calendar",
     "Color-coded calendar shows available, past, closed, and fully booked dates in real time."),
    ("06", "Real-Time Updates via SSE",
     "Server-Sent Events push booking changes to all admin screens instantly — no refresh needed."),
    ("07", "Email Notifications",
     "Automated confirmation, status update, and cancellation emails sent via Gmail SMTP."),
]

C1X = Inches(0.35)
C2X = Inches(7.0)
CW  = Inches(6.0)
RH  = Inches(0.92)
SY  = Inches(1.42)

for i, (num, title, desc) in enumerate(features):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x   = C1X if col == 0 else C2X
    y   = SY + row * RH

    # Number badge
    nb = box(s3, x, y + Inches(0.08), Inches(0.44), Inches(0.44), BURGUNDY)
    tf = nb.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r  = p.add_run()
    r.text = num
    r.font.size = Pt(10)
    r.font.color.rgb = WHITE
    r.font.bold = True
    r.font.name = "Calibri"

    # Title
    txt(s3, title,
        x + Inches(0.56), y, CW - Inches(0.65), Inches(0.38),
        13, BURGUNDY, bold=True)

    # Gold underline for title
    box(s3, x + Inches(0.56), y + Inches(0.36), CW - Inches(0.8), Inches(0.03), LIGHT_GOLD)

    # Description
    txt(s3, desc,
        x + Inches(0.56), y + Inches(0.42), CW - Inches(0.65), Inches(0.48),
        11, BLACK)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — System Architecture (Diagram 1)
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
bg(s4, CREAM)
header(s4, "SYSTEM ARCHITECTURE OVERVIEW")
diagram_placeholder(s4, "Diagram 1  —  System Architecture")

# Small legend strip at bottom
txt(s4,
    "CLIENT LAYER  (Burgundy)   |   APPLICATION LAYER  (Gold)   |   DATA & INTEGRATION LAYER  (Cream)",
    Inches(0.5), Inches(7.1), Inches(12.5), Inches(0.35),
    10, MED_BURG, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Backend Layered Architecture (Diagram 2)
# ─────────────────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
bg(s5, CREAM)
header(s5, "BACKEND LAYERED ARCHITECTURE")
diagram_placeholder(s5, "Diagram 2  —  Backend Layers")

txt(s5,
    "CONTROLLER  (Burgundy)   |   SERVICE  (Medium Burgundy)   |   REPOSITORY  (Gold)   |   ENTITY  (Light Gold)   |   DATABASE  (Cream)",
    Inches(0.5), Inches(7.1), Inches(12.5), Inches(0.35),
    10, MED_BURG, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — End-to-End Booking Data Flow (Diagram 3)
# ─────────────────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
bg(s6, CREAM)
header(s6, "END-TO-END BOOKING DATA FLOW")
diagram_placeholder(s6, "Diagram 3  —  Booking Data Flow")

txt(s6,
    "Customer App  →  Spring Boot API  →  PostgreSQL  →  SSE Service  →  Admin Screen  +  Gmail SMTP",
    Inches(0.5), Inches(7.1), Inches(12.5), Inches(0.35),
    10, MED_BURG, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Thank You
# ─────────────────────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
bg(s7, BURGUNDY)

box(s7, Inches(0), Inches(0),    W, Inches(0.18), GOLD)
box(s7, Inches(0), Inches(7.32), W, Inches(0.18), GOLD)

txt(s7, "Thank You",
    Inches(0), Inches(1.7), W, Inches(1.6),
    68, WHITE, bold=True, align=PP_ALIGN.CENTER, font="Georgia")

box(s7, Inches(4.2), Inches(3.45), Inches(4.9), Inches(0.07), GOLD)

txt(s7, "EventWearify  —  Event Wear Lease System",
    Inches(0), Inches(3.6), W, Inches(0.7),
    22, GOLD, align=PP_ALIGN.CENTER, font="Georgia")

txt(s7, "Jasmine Ciely Palicte",
    Inches(0), Inches(4.55), W, Inches(0.55),
    20, WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(s7, "IT342  —  Systems Integration and Architecture",
    Inches(0), Inches(5.15), W, Inches(0.5),
    15, LIGHT_GOLD, align=PP_ALIGN.CENTER)

txt(s7, "2nd Semester  ·  S.Y. 2025 – 2026",
    Inches(0), Inches(5.7), W, Inches(0.45),
    13, LIGHT_GOLD, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
out = r"c:\2nd Sem 3rd Year\IT342\IT342_EventWearLeaseSystem_G5_Palicte\EventWearify_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
